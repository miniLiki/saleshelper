import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path


@dataclass
class ParsedBlock:
    content: str
    title_path: str = ""
    page_number: int | None = None
    sheet_name: str | None = None


def clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_pdf(data: bytes) -> list[ParsedBlock]:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    blocks: list[ParsedBlock] = []
    for index, page in enumerate(doc, start=1):
        text = clean_text(page.get_text("text"))
        if text:
            blocks.append(ParsedBlock(content=text, title_path=f"第 {index} 页", page_number=index))
    return blocks


def _parse_docx(data: bytes) -> list[ParsedBlock]:
    from docx import Document as DocxDocument

    doc = DocxDocument(BytesIO(data))
    blocks: list[ParsedBlock] = []
    title_stack: list[str] = []
    buffer: list[str] = []
    for paragraph in doc.paragraphs:
        text = clean_text(paragraph.text)
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        if "heading" in style or "标题" in style:
            if buffer:
                blocks.append(ParsedBlock(content="\n".join(buffer), title_path=" / ".join(title_stack)))
                buffer = []
            title_stack = [text]
        else:
            buffer.append(text)
    for table in doc.tables:
        rows = [" | ".join(clean_text(cell.text) for cell in row.cells) for row in table.rows]
        if rows:
            buffer.append("\n".join(rows))
    if buffer:
        blocks.append(ParsedBlock(content="\n".join(buffer), title_path=" / ".join(title_stack)))
    return blocks


def _parse_xlsx(data: bytes) -> list[ParsedBlock]:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(data), data_only=True)
    blocks: list[ParsedBlock] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            blocks.append(
                ParsedBlock(
                    content="\n".join(rows),
                    title_path=f"Sheet: {sheet.title}",
                    sheet_name=sheet.title,
                )
            )
    return blocks


def _parse_pptx(data: bytes) -> list[ParsedBlock]:
    from pptx import Presentation

    presentation = Presentation(BytesIO(data))
    blocks: list[ParsedBlock] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(shape.text)
                if text:
                    texts.append(text)
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if notes:
            texts.append(f"备注：{clean_text(notes)}")
        if texts:
            blocks.append(
                ParsedBlock(content="\n".join(texts), title_path=f"第 {index} 页", page_number=index)
            )
    return blocks


def _parse_text(data: bytes, suffix: str) -> list[ParsedBlock]:
    text = clean_text(data.decode("utf-8", errors="ignore"))
    blocks: list[ParsedBlock] = []
    title_path = ""
    buffer: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if suffix in {".md", ".markdown"} and stripped.startswith("#"):
            if buffer:
                blocks.append(ParsedBlock(content="\n".join(buffer), title_path=title_path))
                buffer = []
            title_path = stripped.lstrip("#").strip()
        elif stripped:
            buffer.append(stripped)
    if buffer:
        blocks.append(ParsedBlock(content="\n".join(buffer), title_path=title_path))
    return blocks or [ParsedBlock(content=text)]


def parse_document_bytes(data: bytes, file_name: str, file_type: str | None = None) -> list[ParsedBlock]:
    suffix = Path(file_name).suffix.lower()
    normalized = (file_type or suffix.lstrip(".")).lower()
    if normalized in {"pdf"} or suffix == ".pdf":
        return _parse_pdf(data)
    if normalized in {"word", "docx"} or suffix == ".docx":
        return _parse_docx(data)
    if normalized in {"excel", "xlsx", "xls"} or suffix in {".xlsx", ".xlsm"}:
        return _parse_xlsx(data)
    if normalized in {"ppt", "pptx"} or suffix == ".pptx":
        return _parse_pptx(data)
    return _parse_text(data, suffix)


def split_into_chunks(blocks: list[ParsedBlock], max_chars: int = 900) -> list[ParsedBlock]:
    chunks: list[ParsedBlock] = []
    for block in blocks:
        paragraphs = [part.strip() for part in re.split(r"\n\s*\n", block.content) if part.strip()]
        current: list[str] = []
        current_len = 0
        for paragraph in paragraphs or [block.content]:
            if current and current_len + len(paragraph) > max_chars:
                chunks.append(
                    ParsedBlock(
                        content=clean_text("\n\n".join(current)),
                        title_path=block.title_path,
                        page_number=block.page_number,
                        sheet_name=block.sheet_name,
                    )
                )
                current = []
                current_len = 0
            if len(paragraph) > max_chars:
                for start in range(0, len(paragraph), max_chars):
                    part = paragraph[start : start + max_chars]
                    chunks.append(
                        ParsedBlock(
                            content=clean_text(part),
                            title_path=block.title_path,
                            page_number=block.page_number,
                            sheet_name=block.sheet_name,
                        )
                    )
            else:
                current.append(paragraph)
                current_len += len(paragraph)
        if current:
            chunks.append(
                ParsedBlock(
                    content=clean_text("\n\n".join(current)),
                    title_path=block.title_path,
                    page_number=block.page_number,
                    sheet_name=block.sheet_name,
                )
            )
    return [chunk for chunk in chunks if chunk.content]
